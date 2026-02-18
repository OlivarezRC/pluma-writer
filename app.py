import streamlit as st
import app.pages as pages
import app.utils as utils
import app.prompts as prompts
import json
import html
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from io import BytesIO

# --- NEW imports ---
from datetime import datetime
from io import BytesIO

# PDF
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import cm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
import os

# --- UI helper: centered "OR" header with lines ---
def or_header(text: str):
    st.markdown(
        """
        <style>
        .or-header{display:flex;align-items:center;gap:.75rem;margin:.5rem 0 1rem;}
        .or-header:before,.or-header:after{content:"";flex:1;border-top:1px solid rgba(128,128,128,.35);}
        .or-header .or-text{font-weight:600;opacity:.85;}
        </style>
        """,
        unsafe_allow_html=True,
    )
    st.markdown(f"<div class='or-header'><span class='or-text'>{text}</span></div>", unsafe_allow_html=True)

# App title
pages.show_home()
pages.show_sidebar()

# Home page
st.header("📝Style Writer")

# # Content input
# st.session_state.content = st.text_area(
#     ":blue[**Content Input for BSP Writing Style:**]", st.session_state.content, 200
# )

# uploaded_files = st.file_uploader(
#     ":blue[**Upload Content Files:**]", 
#     type=["pdf", "docx", "pptx"], 
#     accept_multiple_files=True,
#     help="Upload PDF, Word, or PowerPoint files"
# )
or_header("Input Research Query and Sources for BSP Style Writing")

# Research Query
st.text_input(
    ":blue[**Speech Topic:**]",
    value=st.session_state.get("research_query", ""),
    help="What do you want to research and write about? E.g., 'Use of AI in financial risk management'",
    key="research_query"
)

st.info(
    "Provide at least one source type: **Keywords (for Deep Research)**, **Source Links (URLs)**, or **Upload Attachments**."
)


# --- two-column layout (Col 1 / Col 2) ---
col1, col2 = st.columns([3, 2], gap="small")

with col1:
    # Topics/Keywords for deep research
    st.text_area(
        ":blue[**Keywords (for Deep Research):**]",
        value=st.session_state.get("research_topics", ""),
        height=120,
        help="Space-separated keywords for deep research (e.g., 'AI artificial intelligence machine learning financial portfolio risk prediction')",
        key="research_topics",
    )

    # Dynamic Source Links
    st.markdown(":blue[**Source Links (URLs):**]")
    
    # Initialize source links in session state
    if "source_links" not in st.session_state:
        st.session_state.source_links = [""]
    
    # Display existing links with remove buttons
    links_to_remove = []
    for i, link in enumerate(st.session_state.source_links):
        col_link, col_btn = st.columns([5, 1])
        with col_link:
            new_link = st.text_input(
                f"Link {i+1}",
                value=link,
                key=f"link_input_{i}",
                label_visibility="collapsed",
                placeholder="https://example.com/article"
            )
            st.session_state.source_links[i] = new_link
        with col_btn:
            if st.button("🗑️", key=f"remove_link_{i}", help="Remove this link"):
                links_to_remove.append(i)
    
    # Remove marked links
    for i in reversed(links_to_remove):
        st.session_state.source_links.pop(i)
    
    # Add new link button
    if st.button("➕ Add Source Link", key="add_link"):
        st.session_state.source_links.append("")
        st.rerun()

    # Max iterations for Stage 1 iterative refinement
    st.number_input(
        ":blue[**Research iterations:**]",
        min_value=1,
        max_value=7,
        value=st.session_state.get("max_iterations", 3),
        step=1,
        help="Number of iterative research cycles (1-7). More iterations = more comprehensive but slower.",
        key="max_iterations"
    )

    st.text_area(
        ":blue[**Setting / Location / Conference / Partners (Optional):**]",
        st.session_state.get("context_details", ""),
        height=90,
        help="Optional context to guide greetings and speech considerations.",
        key="context_details",
    )

    MIN_LEN, MAX_LEN, DEFAULT_LEN = 20, 75_000, 1_000

    # One source of truth
    st.session_state.setdefault("max_len", DEFAULT_LEN)
    st.session_state.setdefault("last_updated", None)

    def _update_from_slider():
        st.session_state.last_updated = "slider"
        st.session_state.max_len = st.session_state.max_len_slider

    def _update_from_input():
        st.session_state.last_updated = "input"
        v = st.session_state.max_len_input
        # Coerce + clamp
        try:
            v = int(v)
        except Exception:
            v = MIN_LEN
        st.session_state.max_len = max(MIN_LEN, min(MAX_LEN, v))

    # Keep both widgets synced to the source of truth (avoid ping-pong)
    if st.session_state.last_updated != "slider":
        st.session_state["max_len_slider"] = st.session_state.max_len
    if st.session_state.last_updated != "input":
        st.session_state["max_len_input"] = st.session_state.max_len

    col_slider, col_input = st.columns([3, 1])

    with col_slider:
        st.slider(
            ":blue[**Output Maximum Character Length (75,000 Maximum):**]",
            min_value=MIN_LEN,
            max_value=MAX_LEN,
            key="max_len_slider",
            on_change=_update_from_slider,
            disabled=False,
        )

    with col_input:
        st.number_input(   # use number_input for clean integer UX
            "No. of Characters",
            min_value=MIN_LEN,
            max_value=MAX_LEN,
            step=1,
            key="max_len_input",
            on_change=_update_from_input,
        )

    # Use this in your app
    max_output_length = st.session_state.max_len

with col2:
    
    uploaded_files = st.file_uploader(
        ":blue[**Upload Attachments:**]",
        type=["pdf", "docx", "pptx"],
        accept_multiple_files=True,
        help="Upload PDF, Word, or PowerPoint files as reference materials",
        key="content_upload",
    )

# Extract text from uploaded files as attachments
attachment_contents = []
if uploaded_files:
    for uploaded_file in uploaded_files:
        file_type = uploaded_file.name.split('.')[-1].lower()
        raw_bytes = uploaded_file.getvalue()
        file_text = ""
        
        if file_type == 'pdf':
            pdf_reader = PdfReader(BytesIO(raw_bytes))
            for page in pdf_reader.pages:
                file_text += page.extract_text() + "\n"
                
        elif file_type == 'docx':
            doc = Document(BytesIO(raw_bytes))
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    file_text += paragraph.text + "\n"
                    
        elif file_type == 'pptx':
            prs = Presentation(BytesIO(raw_bytes))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            file_text += shape.text + "\n"
        
        if file_text.strip() or raw_bytes:
            attachment_contents.append({
                "filename": uploaded_file.name,
                "file_type": file_type,
                "content": file_text,
                "bytes": raw_bytes
            })

# Show preview of attachments
if attachment_contents:
    with st.expander(f"📎 Attached Files: {len(attachment_contents)} file(s)", expanded=False):
        for att in attachment_contents:
            st.markdown(f"**{att['filename']}** ({len(att['content'])} characters)")
            st.text_area(
                f"Preview: {att['filename']}", 
                att['content'][:500] + ("..." if len(att['content']) > 500 else ""),
                height=150,
                key=f"preview_{att['filename']}",
                disabled=True
            )

# Source readiness block
query_filled = bool(st.session_state.get("research_query", "").strip())
topics_filled = bool(st.session_state.get("research_topics", "").strip())
links_filled = any(link.strip() for link in st.session_state.get("source_links", []))
attachments_filled = len(attachment_contents) > 0
has_any_source = topics_filled or links_filled or attachments_filled

if not query_filled and not has_any_source:
    readiness_footer = "You haven't filled either of the three source types and Research Query / Speech Topic"
elif not query_filled:
    readiness_footer = "You still don't have a Research Query / Speech Topic:"
elif not has_any_source:
    readiness_footer = "You still haven't filled either of the three source types"
else:
    readiness_footer = "✅ Requirements satisfied."

readiness_text = (
    ("You are good to go!" if (query_filled and has_any_source) else "Either add a Topics / Keywords to be used for deep research, source links to be searched, or upload attachments to be extracted.") + "\n\n"
    f"- {'✅' if query_filled else '⬜'} Research Query / Speech Topic\n"
    f"- {'✅' if topics_filled else '⬜'} Topics / Keywords\n"
    f"- {'✅' if links_filled else '⬜'} Source Links\n"
    f"- {'✅' if attachments_filled else '⬜'} Upload Attachments\n\n"
    f"{readiness_footer}"
)

if query_filled and has_any_source:
    st.success(readiness_text)
elif not query_filled:
    st.info(readiness_text)
else:
    st.warning(readiness_text)

# Extracting the styles and creating combined display options
styles_data = utils.get_styles()
def _to_kv_rows(data):
    if not data:
        return []
    if isinstance(data, str):
        return [{"Key": "text", "Value": data}]
    rows = []
    if isinstance(data, dict):
        for key, value in data.items():
            if isinstance(value, (dict, list)):
                value = json.dumps(value, ensure_ascii=False, indent=2)
            rows.append({"Key": str(key), "Value": value})
    return rows


def _normalize_style_rules(style_rules):
    if not style_rules:
        return {}
    if isinstance(style_rules, str):
        try:
            return json.loads(style_rules)
        except Exception:
            return {"style_rules": style_rules}
    if isinstance(style_rules, dict):
        return style_rules
    return {"style_rules": str(style_rules)}


def _schema_to_values(schema):
    if not isinstance(schema, dict):
        return schema
    if "enum" in schema and isinstance(schema["enum"], list):
        return schema["enum"][0] if schema["enum"] else None
    if schema.get("type") == "array" and isinstance(schema.get("items"), dict):
        items = schema["items"]
        if "enum" in items and isinstance(items["enum"], list):
            return items["enum"]
        return []
    if "properties" in schema and isinstance(schema["properties"], dict):
        return {k: _schema_to_values(v) for k, v in schema["properties"].items()}
    if "items" in schema and isinstance(schema["items"], dict):
        return _schema_to_values(schema["items"])
    return schema


def _extract_style_rules(item):
    rules = item.get("style") or item.get("style_instructions")
    rules = _normalize_style_rules(rules)
    if rules:
        return rules
    # Check nested properties payload
    props = item.get("properties") or {}
    if isinstance(props, str):
        try:
            props = json.loads(props)
        except Exception:
            props = {}
    if isinstance(props, dict):
        nested = props.get("style_instructions") or props.get("style")
        if isinstance(nested, dict):
            nested_values = _schema_to_values(nested)
            if nested_values:
                return nested_values
        nested = _normalize_style_rules(nested)
        if nested:
            return nested
    # Fallback: build from known top-level fields
    fallback_keys = [
        "register_profile",
        "stylistics",
        "pragmatics",
        "semantics_frames",
        "discourse_structure",
        "signature_moves",
        "evidence_spans",
    ]
    compiled = {}
    for key in fallback_keys:
        if key in item and item[key]:
            compiled[key] = item[key]
    if not compiled and isinstance(props, dict):
        for key in fallback_keys:
            if key in props and props[key]:
                compiled[key] = _schema_to_values(props[key])
    return compiled


def _render_rule_section(title, value):
    if value is None or value == "":
        st.caption("No data.")
        return
    if isinstance(value, list):
        items = [str(v) for v in value if v is not None and str(v).strip() != ""]
        if items:
            st.markdown(
                """
                <style>
                .rule-list{margin:0.25rem 0 0.5rem 1.1rem;}
                .rule-list li{margin:0.2rem 0;}
                </style>
                """,
                unsafe_allow_html=True,
            )
            st.markdown(
                "<ul class='rule-list'>" + "".join([f"<li>{i}</li>" for i in items]) + "</ul>",
                unsafe_allow_html=True,
            )
        else:
            st.caption("No data.")
        return
    if isinstance(value, dict):
        rows = []
        complex_keys = []
        for k, v in value.items():
            if isinstance(v, (dict, list)):
                complex_keys.append((k, v))
            else:
                rows.append({"Key": str(k), "Value": v})
        if rows:
            st.markdown(
                """
                <style>
                .rule-kv{display:flex;flex-direction:column;gap:0.35rem;margin:0.35rem 0 0.5rem;}
                .rule-kv .rule-item{display:flex;flex-wrap:wrap;gap:0.35rem;}
                .rule-kv .rule-key{font-weight:600;}
                </style>
                """,
                unsafe_allow_html=True,
            )
            items_html = "".join(
                [
                    "<div class='rule-item'><span class='rule-key'>"
                    + str(r["Key"])
                    + ":</span><span class='rule-val'>"
                    + str(r["Value"])
                    + "</span></div>"
                    for r in rows
                ]
            )
            st.markdown(f"<div class='rule-kv'>{items_html}</div>", unsafe_allow_html=True)
        for k, v in complex_keys:
            with st.expander(str(k), expanded=False):
                _render_rule_section(str(k), v)
        return
    st.markdown(
        """
        <style>
        .rule-single{margin:0.35rem 0 0.5rem;}
        </style>
        """,
        unsafe_allow_html=True,
    )
    st.markdown(f"<div class='rule-single'>{value}</div>", unsafe_allow_html=True)


def _format_section_title(text: str) -> str:
    return str(text).replace("_", " ").replace("-", " ").strip()


def _filter_global_rulebook(rulebook: dict) -> dict:
    if not isinstance(rulebook, dict):
        return {}
    excluded = {
        "id",
        "_rid",
        "_self",
        "_etag",
        "_attachments",
        "_ts",
        "PartitionKey",
        "doc_kind",
        "scope",
        "version",
        "created_at",
        "coverage_metrics",
        "confidence",
        "evidence_spans",
        "speech_ids",
        "speech_count",
    }
    return {k: v for k, v in rulebook.items() if k not in excluded}


def _coerce_global_rulebook():
    rulebook = st.session_state.get("global_rulebook")
    if rulebook:
        return rulebook
    global_rules = st.session_state.get("global_rules")
    if isinstance(global_rules, dict):
        return global_rules
    if isinstance(global_rules, str) and global_rules.strip():
        try:
            return json.loads(global_rules)
        except Exception:
            return {"global_rules": global_rules}
    return {}


def _split_style_name(name):
    if "_" in name:
        speaker, audience = name.split("_", 1)
        return speaker.strip(), audience.strip()
    return name.strip(), "General"

speakers = []
audience_by_speaker = {}
for item in styles_data:
    name = item.get("name") or item.get("container_key") or item.get("id")
    if not name:
        speaker = item.get("speaker")
        audience = item.get("audience_setting_classification")
        if speaker and audience:
            name = f"{speaker}_{audience}"
        else:
            continue
    speaker, audience = _split_style_name(str(name))
    if speaker not in audience_by_speaker:
        audience_by_speaker[speaker] = []
        speakers.append(speaker)
    if audience not in audience_by_speaker[speaker]:
        audience_by_speaker[speaker].append(audience)

selected_speaker = st.selectbox(":blue[**Select Speaker:**]", options=speakers, index=None, key="selected_speaker_input")
selected_audience = st.selectbox(
    ":blue[**Select Audience/Setting:**]",
    options=audience_by_speaker.get(selected_speaker, []) if selected_speaker else [],
    index=None,
    key="selected_audience_input"
)

if not speakers:
    st.warning("No speakers available. Check the styles container data.")

# Assigning the selected style to the session state
if selected_speaker and selected_audience:
    selected_style_name = f"{selected_speaker}_{selected_audience}"
    filtered = utils.get_style_fingerprint(selected_speaker, selected_audience)
    if not filtered:
        filtered = next(
            (
                item for item in styles_data
                if str(item.get("name") or item.get("container_key") or item.get("id") or "")
                == selected_style_name
                or (
                    item.get("speaker") == selected_speaker
                    and item.get("audience_setting_classification") == selected_audience
                )
            ),
            None
        )

    st.session_state.selected_style_doc = filtered or {}

    if filtered:
        style_rules = _extract_style_rules(filtered)
        st.session_state.style_rules = style_rules
        st.session_state.style = json.dumps(style_rules, ensure_ascii=False, indent=2) if style_rules else ""
        example = filtered.get("example")
        if not example:
            evidence = filtered.get("evidence_spans") or []
            example = "\n".join([e for e in evidence if e]) if evidence else ""
        st.session_state.example = example
        st.session_state.styleId = selected_style_name
        rulebook_id = filtered.get("global_rulebook_id")
        rulebook = utils.get_rulebook(rulebook_id)
        st.session_state.global_rulebook = rulebook or {}
        st.session_state.global_rules = utils.extract_rulebook_text(rulebook)

show_loaded_rules = False

if show_loaded_rules:
    with st.container(border=True):
        st.subheader("Loaded Style & Global Rules")
        if not (selected_speaker and selected_audience):
            st.info("Select a speaker and audience/setting to view the loaded rules.")
        else:
            selected_doc = st.session_state.get("selected_style_doc") or {}
            style_source = st.session_state.get("style_rules") or _extract_style_rules(selected_doc) or st.session_state.get("style") or {}
            global_rulebook = _coerce_global_rulebook()

            st.markdown("**Style Rules**")
            if style_source:
                if isinstance(style_source, dict):
                    rule_items = list(style_source.items())
                    mid_point = (len(rule_items) + 1) // 2
                    col1, col2 = st.columns(2)
                    with col1:
                        for section, value in rule_items[:mid_point]:
                            section_title = _format_section_title(section)
                            with st.expander(section_title, expanded=False):
                                _render_rule_section(section_title, value)
                    with col2:
                        for section, value in rule_items[mid_point:]:
                            section_title = _format_section_title(section)
                            with st.expander(section_title, expanded=False):
                                _render_rule_section(section_title, value)
                else:
                    _render_rule_section("Style Rules", style_source)
            else:
                st.caption("No style rules loaded.")

            st.markdown("**Global Rules**")
            if global_rulebook:
                if isinstance(global_rulebook, dict):
                    filtered_rulebook = _filter_global_rulebook(global_rulebook)
                    rule_items = list(filtered_rulebook.items())
                    if not rule_items:
                        st.caption("No global rules loaded.")
                        rule_items = []
                    mid_point = (len(rule_items) + 1) // 2
                    col1, col2 = st.columns(2)
                    with col1:
                        for section, value in rule_items[:mid_point]:
                            section_title = _format_section_title(section)
                            with st.expander(section_title, expanded=False):
                                _render_rule_section(section_title, value)
                    with col2:
                        for section, value in rule_items[mid_point:]:
                            section_title = _format_section_title(section)
                            with st.expander(section_title, expanded=False):
                                _render_rule_section(section_title, value)
                else:
                    _render_rule_section("Global Rules", global_rulebook)
            else:
                st.caption("No global rules loaded.")

            with st.expander("Style Debug", expanded=False):
                st.write({
                    "selected_speaker": selected_speaker,
                    "selected_audience": selected_audience,
                    "style_doc_loaded": bool(selected_doc),
                    "style_doc_keys": list(selected_doc.keys())[:25],
                    "has_style_instructions": "style_instructions" in selected_doc,
                    "has_register_profile": "register_profile" in selected_doc,
                    "style_rules_keys": list(style_source.keys())[:25] if isinstance(style_source, dict) else [],
                })
        
# st.session_state.style = st.text_area(":blue[**Style:**]", st.session_state.style)

# Show the example style
guidelines = st.session_state.locals.get("relevant_guidelines", {})
guidelines_summary = st.session_state.locals.get("guideline_summaries", {}) 
selected_guidelines = []

st.write(":blue[**Select Editorial Style Guides:**]")

# Tooltip for guideline summary in the UI
def render_guideline_checkbox(section_name: str, content: str, col_key_prefix: str):
    default_checked = section_name in ["COMMON GRAMMATICAL ERRORS", "WRITING LETTERS"]
    tooltip = guidelines_summary.get(section_name, None)  # one-sentence summary for hover
    if st.checkbox(
        section_name,
        value=default_checked,
        key=f"{col_key_prefix}_{section_name}",
        help=tooltip  # <-- hover tooltip appears on the ⓘ icon and on hover
    ):
        selected_guidelines.append(content)

# Create a checkbox for each guideline section
if guidelines:
    with st.container(border=True):
        # Create two columns
        col1, col2 = st.columns(2)

        # Split guidelines into two halves
        guideline_items = list(guidelines.items())
        mid_point = len(guideline_items) // 2

        # First column
        with col1:
            for section_name, content in guideline_items[:mid_point]:
                render_guideline_checkbox(section_name, content, "col1")
                # default = section_name in ["COMMON GRAMMATICAL ERRORS", "WRITING LETTERS"]
                # if st.checkbox(section_name, value=default, key=f"col1_{section_name}"):
                #     selected_guidelines.append(content)

        # Second column
        with col2:
            for section_name, content in guideline_items[mid_point:]:
                render_guideline_checkbox(section_name, content, "col2")
                # default = section_name in ["COMMON GRAMMATICAL ERRORS", "WRITING LETTERS"]
                # if st.checkbox(section_name, value=default, key=f"col2_{section_name}"):
                #     selected_guidelines.append(content)
else:
    st.warning("No guidelines available in the local data.")

# Join all selected guidelines with newlines and store in session state
st.session_state.guidelines = "\n".join(selected_guidelines)

# Show the combined guidelines in a text area
# st.text_area(":blue[**Relevant Guidelines:**]", st.session_state.guidelines, height=200)


# --- NEW helpers: make DOCX/PDF from text ---

def make_docx_bytes(text: str, title: str | None = None) -> bytes:
    """Return a .docx file (bytes) with a title and body paragraphs."""
    doc = Document()
    if title:
        doc.add_heading(title, level=1)
    # Split into paragraphs on blank lines while preserving line breaks
    for block in text.replace("\r\n", "\n").split("\n\n"):
        p = doc.add_paragraph()
        for line in block.split("\n"):
            if line.strip():
                p.add_run(line)
            p.add_run("\n")
    bio = BytesIO()
    doc.save(bio)
    return bio.getvalue()


def _register_pdf_font_if_available():
    """Optionally register DejaVuSans for better Unicode PDF rendering."""
    try:
        font_path = os.path.join("assets", "DejaVuSans.ttf")
        if os.path.exists(font_path):
            pdfmetrics.registerFont(TTFont("DejaVuSans", font_path))
            return "DejaVuSans"
    except Exception:
        pass
    # Fallback to built-in Helvetica (ASCII/Latin-1 safe)
    return "Helvetica"


def make_pdf_bytes(text: str, title: str | None = None) -> bytes:
    """Return a PDF (bytes) using ReportLab."""
    font_name = _register_pdf_font_if_available()

    buf = BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=A4,
        leftMargin=2 * cm,
        rightMargin=2 * cm,
        topMargin=2 * cm,
        bottomMargin=2 * cm,
        title=title or "Rewrite",
        author="Style Writer",
    )

    styles = getSampleStyleSheet()
    base = styles["BodyText"]
    base.fontName = font_name
    base.fontSize = 11
    base.leading = 14

    title_style = ParagraphStyle(
        "Title",
        parent=styles["Heading1"],
        fontName=font_name,
        spaceAfter=12,
    )

    story = []
    if title:
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 8))

    # Turn double newlines into paragraph breaks; single newlines stay inside a paragraph
    for block in text.replace("\r\n", "\n").split("\n\n"):
        # Escape simple HTML-sensitive chars for Paragraph
        block = block.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        block = block.replace("\n", "<br/>")
        story.append(Paragraph(block, base))
        story.append(Spacer(1, 6))

    doc.build(story)
    return buf.getvalue()


if st.button(
    ":blue[**Generate Speech with Complete Pipeline**]",
    key="generate_speech",
    disabled=(
        not st.session_state.get("research_query", "").strip()
        or (
            not st.session_state.get("research_topics", "").strip()
            and (not st.session_state.source_links or all(not link.strip() for link in st.session_state.source_links))
            and not attachment_contents
        )
    ),
    help="Requires: Research Query plus at least one source (Topics/Keywords, Source Link, or Attachment)"
):
    # Import complete pipeline
    import asyncio
    from app.writer_main import process_with_iterative_refinement_and_style
    import sys
    from io import StringIO
    import re
    
    with st.container(border=True):
        # Prepare sources dictionary
        user_query = st.session_state.get("research_query", "")
        user_topics = st.session_state.get("research_topics", "")
        context_details = st.session_state.get("context_details", "")
        
        # Filter out empty links
        source_links = [link.strip() for link in st.session_state.source_links if link.strip()]
        
        # Prepare attachments (structured so backend can use Document Intelligence + LLM metadata extraction)
        attachments = attachment_contents
        
        sources = {
            "topics": user_topics,
            "links": source_links,
            "attachments": attachments
        }
        
        # Prepare style from selected speaker/audience
        selected_style_doc = st.session_state.get("selected_style_doc") or {}
        selected_style = dict(selected_style_doc) if isinstance(selected_style_doc, dict) else {}

        if st.session_state.get("style"):
            selected_style["style_description"] = st.session_state.get("style")
        if st.session_state.get("global_rules"):
            selected_style["global_rules"] = st.session_state.get("global_rules")
        if st.session_state.get("guidelines"):
            selected_style["guidelines"] = st.session_state.get("guidelines")
        if st.session_state.get("example"):
            selected_style["example"] = st.session_state.get("example")

        # Responsive stage cards: background worker + structured event queue
        import threading
        import queue
        import time

        stage_names = {
            1: "Iterative Refinement",
            2: "Retrieving Writing Style",
            3: "Generating Styled Output",
            4: "Verifying Citations",
            5: "Converting to APA Format",
            6: "Plagiarism Detection",
            7: "BSP Policy Alignment Check"
        }

        stage_store = {
            i: {
                "status": "pending",
                "lines": [],
                "metrics": {},
                "style_preview": "",
                "initial_output_preview": "",
                "generated_output": "",
            }
            for i in range(1, 8)
        }

        stage_placeholders = {i: st.empty() for i in range(1, 8)}
        event_queue = queue.Queue()
        result_holder = {"results": None, "error": None}

        def render_stage(stage_num: int):
            state = stage_store[stage_num]
            icon_map = {
                "pending": "⏳",
                "running": "🔄",
                "complete": "✅",
                "error": "❌",
            }
            icon = icon_map.get(state["status"], "⏳")
            suffix = " - Complete" if state["status"] == "complete" else ""

            with stage_placeholders[stage_num].container(border=True):
                st.markdown(f"### {icon} Stage {stage_num}: {stage_names[stage_num]}{suffix}")
                st.caption(
                    "Running live" if state["status"] == "running"
                    else "Completed" if state["status"] == "complete"
                    else "Waiting" if state["status"] == "pending"
                    else "Needs attention"
                )

                if state["lines"]:
                    st.code("\n".join(state["lines"][-8:]), language=None)

                if state["metrics"]:
                    cols = st.columns(len(state["metrics"]))
                    for idx, (key, value) in enumerate(state["metrics"].items()):
                        with cols[idx]:
                            st.metric(key, value)

                if stage_num == 2 and state["style_preview"]:
                    st.markdown("**Style Preview**")
                    st.code(state["style_preview"], language=None)

                if stage_num == 3 and state["initial_output_preview"]:
                    st.markdown("**Initial Output Preview**")
                    st.code(state["initial_output_preview"], language=None)

                if stage_num == 3 and state["generated_output"]:
                    st.markdown("**Generated Styled Output**")
                    st.code(state["generated_output"], language=None)

        def render_all_stages():
            for i in range(1, 8):
                render_stage(i)

        # Capture selected iteration count on the main thread (thread-safe)
        selected_max_iterations = int(st.session_state.get("max_iterations", 3))
        selected_max_iterations = max(1, min(7, selected_max_iterations))
        selected_max_output_length = int(st.session_state.get("max_len", 1000))
        selected_max_output_length = max(20, min(75000, selected_max_output_length))
        stage_store[1]["lines"].append(f"Configured iterations: {selected_max_iterations}")

        def worker():
            try:
                import builtins
                current_stage = 0

                def emit_text(stage_num: int, text: str):
                    event_queue.put({"type": "stage_text", "stage": stage_num, "text": text})

                def parse_print_line(text: str):
                    nonlocal current_stage
                    t = text.strip()
                    if not t:
                        return

                    if "STEP 1:" in t or "ITERATIVE REFINEMENT" in t:
                        current_stage = 1
                    elif "STEP 2:" in t or "RETRIEVING WRITING STYLE" in t:
                        current_stage = 2
                    elif "STEP 3:" in t or "GENERATING STYLED OUTPUT" in t:
                        current_stage = 3
                    elif "STEP 4:" in t or "VERIFYING STYLED OUTPUT CITATIONS" in t:
                        current_stage = 4
                    elif "STEP 5:" in t or "CONVERTING TO APA FORMAT" in t:
                        current_stage = 5
                    elif "STEP 6:" in t or "PLAGIARISM DETECTION" in t:
                        current_stage = 6
                    elif "STEP 7:" in t or "BSP POLICY ALIGNMENT CHECK" in t:
                        current_stage = 7

                    # Stage 1 (evidence + iteration flow)
                    if current_stage == 1:
                        if any(k in t for k in [
                            "Evidence IDs starting at",
                            "Extracted",
                            "Link processing:",
                            "Generating critique",
                            "Generating adjustments",
                            "Next iteration will use refined query",
                            "Total evidence collected",
                        ]):
                            emit_text(1, t)

                    # Stage 2 (style retrieval metadata)
                    if current_stage == 2:
                        if any(k in t for k in ["Retrieved style:", "Speaker:", "Audience:", "No style found"]):
                            emit_text(2, t)

                    # Stage 3 (sample speeches + initial output)
                    if current_stage == 3:
                        if any(k in t for k in [
                            "Fetching sample speeches",
                            "Fetched",
                            "Speech 1:",
                            "Speech 2:",
                            "Speech 3:",
                            "Styled output generated successfully",
                            "Style applied:",
                            "Output length:",
                            "All citations valid",
                        ]):
                            emit_text(3, t)
                        if "Content received, first 100 chars:" in t:
                            preview = t.split("Content received, first 100 chars:", 1)[-1].strip()
                            event_queue.put({"type": "stage_preview", "stage": 3, "text": preview})

                    # Stage 4 (verification summary + autofix)
                    if current_stage == 4:
                        if any(k in t for k in [
                            "VERIFYING STYLED OUTPUT CITATIONS",
                            "Evidence store:",
                            "Parsed",
                            "Verifying segments",
                            "Verified ",
                            "Verification complete",
                            "Verification rate",
                            "Sample unverified segments",
                            "[AUTO-FIX]",
                        ]):
                            emit_text(4, t)

                    # Stage 5 (APA conversion)
                    if current_stage == 5:
                        if any(k in t for k in [
                            "CONVERTING TO APA FORMAT",
                            "Evidence store:",
                            "Building APA citations",
                            "Built citation map",
                            "Unique sources",
                            "Converting [ENN] citations",
                            "Generating APA",
                            "APA conversion complete",
                            "Citations converted",
                            "References generated",
                        ]):
                            emit_text(5, t)

                    # Stage 6 (concise plagiarism flow, skip noisy batches)
                    if current_stage == 6:
                        if any(skip in t for skip in ["Batch ", "Source "]):
                            return
                        if any(k in t for k in [
                            "PLAGIARISM DETECTION ANALYSIS",
                            "[Step 0]",
                            "[Step 1]",
                            "[Step 2]",
                            "[Step 3]",
                            "[Step 4]",
                            "[Step 5]",
                            "[Step 6]",
                            "Created ",
                            "Collected ",
                            "potential source",
                            "Overall risk level",
                            "SUMMARY:",
                            "Overall risk:",
                            "High risk:",
                            "Medium risk:",
                            "Low risk:",
                            "Clean:",
                        ]):
                            emit_text(6, t)

                    # Stage 7 (policy summary)
                    if current_stage == 7:
                        if any(k in t for k in [
                            "BSP POLICY ALIGNMENT CHECK",
                            "Overall Compliance:",
                            "Compliance Score:",
                            "Violations Found:",
                            "Commendations:",
                            "BSP Circulars Referenced:",
                            "RECOMMENDATION:",
                            "Policy alignment check complete",
                        ]):
                            emit_text(7, t)

                original_print = builtins.print

                def custom_print(*args, **kwargs):
                    text = " ".join(str(a) for a in args)
                    parse_print_line(text)
                    original_print(*args, **kwargs)

                builtins.print = custom_print

                results = asyncio.run(
                    process_with_iterative_refinement_and_style(
                        query=user_query,
                        sources=sources,
                        max_iterations=selected_max_iterations,
                        max_output_length=selected_max_output_length,
                        context_details=context_details,
                        style=selected_style if selected_style else None,
                        enable_policy_check=True,
                        progress_callback=lambda event: event_queue.put(event),
                    )
                )
                builtins.print = original_print
                event_queue.put({"type": "pipeline_complete", "results": results})
            except Exception as exc:
                try:
                    builtins.print = original_print
                except Exception:
                    pass
                event_queue.put({"type": "pipeline_error", "error": str(exc)})

        st.markdown("### 🔄 Writer Pipeline Running...")
        st.markdown("---")
        render_all_stages()

        pipeline_thread = threading.Thread(target=worker, daemon=True)
        pipeline_thread.start()

        while pipeline_thread.is_alive() or not event_queue.empty():
            updated = False
            while not event_queue.empty():
                event = event_queue.get()
                updated = True
                event_type = event.get("type")

                if event_type == "stage_started":
                    stage = event.get("stage")
                    if stage in stage_store:
                        stage_store[stage]["status"] = "running"
                elif event_type == "stage_done":
                    stage = event.get("stage")
                    if stage in stage_store:
                        stage_store[stage]["status"] = "complete"
                elif event_type == "stage_text":
                    stage = event.get("stage")
                    text = str(event.get("text", "")).strip()
                    if stage in stage_store and text and text not in stage_store[stage]["lines"]:
                        stage_store[stage]["lines"].append(text)
                        if len(stage_store[stage]["lines"]) > 80:
                            stage_store[stage]["lines"] = stage_store[stage]["lines"][-80:]
                elif event_type == "stage_preview":
                    stage = event.get("stage")
                    text = str(event.get("text", "")).strip()
                    if stage in stage_store and text:
                        stage_store[stage]["initial_output_preview"] = text
                elif event_type == "stage_metric":
                    stage = event.get("stage")
                    key = str(event.get("key", "")).strip()
                    value = event.get("value", "")
                    if stage in stage_store and key:
                        stage_store[stage]["metrics"][key] = value
                elif event_type == "pipeline_complete":
                    results = event.get("results") or {}
                    result_holder["results"] = results
                    for i in range(1, 8):
                        if stage_store[i]["status"] != "error":
                            stage_store[i]["status"] = "complete"

                    style_used = results.get("style_used", {}) if isinstance(results, dict) else {}
                    style_name = style_used.get("name", "Unknown") if isinstance(style_used, dict) else "Unknown"
                    speaker = style_used.get("speaker", "Unknown") if isinstance(style_used, dict) else "Unknown"
                    audience = style_used.get("audience", "General") if isinstance(style_used, dict) else "General"

                    preview = ""
                    if selected_style and isinstance(selected_style, dict):
                        for key in ["transcript_text", "style_description", "style", "guidelines", "example"]:
                            value = selected_style.get(key, "")
                            if isinstance(value, str) and value.strip():
                                preview = value.strip()[:700]
                                break

                    stage_store[2]["lines"] = [
                        f"Style: {style_name}",
                        f"Speaker: {speaker}",
                        f"Audience: {audience}",
                    ]
                    if selected_style and isinstance(selected_style, dict):
                        style_id = selected_style.get("id") or selected_style.get("style_id") or selected_style.get("speaker_style_id")
                        if style_id:
                            stage_store[2]["lines"].append(f"Style ID: {style_id}")
                    stage_store[2]["style_preview"] = preview

                    styled_payload = results.get("styled_output") if isinstance(results, dict) else {}
                    styled_output = ""
                    if isinstance(styled_payload, dict):
                        styled_output = styled_payload.get("styled_output", "")
                    if not styled_output and isinstance(results.get("styled_output_apa"), dict):
                        styled_output = results.get("styled_output_apa", {}).get("apa_output", "")
                    if not styled_output and isinstance(results.get("final_summary"), dict):
                        styled_output = results.get("final_summary", {}).get("summary", "")

                    stage_store[3]["generated_output"] = styled_output
                    stage_store[3]["metrics"] = {
                        "Output Length": len(styled_output) if styled_output else 0,
                        "Citations": (styled_payload or {}).get("citations_found", 0) if isinstance(styled_payload, dict) else 0,
                        "Evidence IDs": (styled_payload or {}).get("unique_evidence_cited", 0) if isinstance(styled_payload, dict) else 0,
                    }
                    if styled_output and not stage_store[3]["initial_output_preview"]:
                        stage_store[3]["initial_output_preview"] = styled_output[:280]
                elif event_type == "pipeline_error":
                    result_holder["error"] = event.get("error", "Unknown pipeline error")
                    # Mark current incomplete stages as error
                    for i in range(1, 8):
                        if stage_store[i]["status"] in ["pending", "running"]:
                            stage_store[i]["status"] = "error"

            if updated:
                render_all_stages()

            time.sleep(0.15)

        render_all_stages()

        if result_holder["error"]:
            st.error(f"Pipeline failed: {result_holder['error']}")
            st.stop()

        results = result_holder["results"]
        if not results:
            st.error("Pipeline returned no results.")
            st.stop()

        st.session_state["pipeline_results"] = results
        st.success("Pipeline completed with responsive stage updates.")

        st.markdown("---")
        st.markdown("### 📝 Final Output")

        # Prefer APA output, then styled output, then raw summary
        apa_result = results.get("styled_output_apa", {}) if isinstance(results, dict) else {}
        styled_result = results.get("styled_output", {}) if isinstance(results, dict) else {}
        final_summary = results.get("final_summary", {}) if isinstance(results, dict) else {}

        if isinstance(apa_result, dict) and apa_result.get("success") and apa_result.get("apa_output"):
            final_output_text = apa_result.get("apa_output", "")
            output_label = "APA Final Speech"
        elif isinstance(styled_result, dict) and styled_result.get("success") and styled_result.get("styled_output"):
            final_output_text = styled_result.get("styled_output", "")
            output_label = "Styled Final Speech"
        else:
            final_output_text = final_summary.get("summary", "") if isinstance(final_summary, dict) else ""
            output_label = "Final Summary (Fallback)"

        if final_output_text:
            st.text_area(
                output_label,
                final_output_text,
                height=420,
                key="final_output_display",
            )

            ts = datetime.now().strftime("%Y%m%d-%H%M%S")
            speaker = st.session_state.get("selected_speaker_input", "Speaker")
            base_name = f"speech_{speaker}_{ts}"

            col_docx, col_pdf = st.columns(2)
            with col_docx:
                docx_bytes = make_docx_bytes(final_output_text, title=f"Speech - {speaker}")
                st.download_button(
                    "⬇️ Download DOCX",
                    data=docx_bytes,
                    file_name=f"{base_name}.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    use_container_width=True,
                    key="download_final_docx"
                )
            with col_pdf:
                pdf_bytes = make_pdf_bytes(final_output_text, title=f"Speech - {speaker}")
                st.download_button(
                    "⬇️ Download PDF",
                    data=pdf_bytes,
                    file_name=f"{base_name}.pdf",
                    mime="application/pdf",
                    use_container_width=True,
                    key="download_final_pdf"
                )
        else:
            st.warning("No final output text was produced by the pipeline.")

        st.stop()
        
        # Create stage placeholders
        st.markdown("### 🔄 Writer Pipeline Running...")
        st.markdown("---")
        
        # Create status containers for live updates (all stages use st.status for consistency)
        stage_1_status = st.status("⏳ Stage 1: Iterative Refinement", expanded=True)
        stage_2_status = st.status("⏳ Stage 2: Retrieving Writing Style", expanded=False)
        stage_3_status = st.status("⏳ Stage 3: Generating Styled Output", expanded=False)
        stage_4_status = st.status("⏳ Stage 4: Verifying Citations", expanded=False)
        stage_5_status = st.status("⏳ Stage 5: Converting to APA Format", expanded=False)
        stage_6_status = st.status("⏳ Stage 6: Plagiarism Detection", expanded=False)
        stage_7_status = st.status("⏳ Stage 7: BSP Policy Alignment Check", expanded=False)
        
        stage_containers = {
            1: stage_1_status,
            2: stage_2_status,
            3: stage_3_status,
            4: stage_4_status,
            5: stage_5_status,
            6: stage_6_status,
            7: stage_7_status
        }
        stage_status = {i: "pending" for i in range(1, 8)}
        
        # Stage names
        stage_names = {
            1: "Iterative Refinement",
            2: "Retrieving Writing Style",
            3: "Generating Styled Output",
            4: "Verifying Citations",
            5: "Converting to APA Format",
            6: "Plagiarism Detection",
            7: "BSP Policy Alignment Check"
        }
        
        # Function to update stage card
        def update_stage_card(stage_num, status, title, content="", metrics=None):
            """
            status: 'running', 'complete', 'error', 'pending'
            """
            stage_status[stage_num] = status
            
            # Icon based on status
            icons = {
                'pending': '⏳',
                'running': '🔄',
                'complete': '✅',
                'error': '❌'
            }
            
            icon = icons.get(status, '⏳')
            status_container = stage_containers[stage_num]
            
            # Update the status container label and state
            if status == 'complete':
                status_container.update(label=f"✅ Stage {stage_num}: {title} - Complete", state="complete", expanded=False)
            elif status == 'running':
                status_container.update(label=f"🔄 Stage {stage_num}: {title}", state="running", expanded=False)
            elif status == 'error':
                status_container.update(label=f"❌ Stage {stage_num}: {title} - Error", state="error", expanded=True)
            
            # Add content inside the status container
            if content or metrics:
                with status_container:
                    if content:
                        st.text(content)
                    if metrics:
                        cols = st.columns(len(metrics))
                        for idx, (label, value) in enumerate(metrics.items()):
                            with cols[idx]:
                                st.metric(label, value)
        
        # Track messages per stage for detailed display
        stage_messages = {i: [] for i in range(1, 8)}
        
        # Capture stdout to parse progress
        class ProgressCapture:
            def __init__(self):
                self.current_stage = 0
                self.buffer = []
                self.original_stdout = sys.stdout
                self.current_iteration = 0
                self.iteration_data = {}
                self.last_rendered_iteration = 0
                
            def parse_and_update(self, text):
                # Store message for current stage
                if self.current_stage > 0:
                    stage_messages[self.current_stage].append(text)
                
                # Check for stage markers
                if "STEP 1:" in text or "ITERATIVE REFINEMENT" in text:
                    self.current_stage = 1
                    self.update_stage_1_card('running')
                elif "STEP 2:" in text or "RETRIEVING WRITING STYLE" in text:
                    if self.current_stage == 1:
                        self.update_stage_1_card('complete')
                    self.current_stage = 2
                    update_stage_card(2, 'running', stage_names[2])
                elif "STEP 3:" in text or "GENERATING STYLED OUTPUT" in text:
                    if self.current_stage == 2:
                        update_stage_card(2, 'complete', stage_names[2])
                    self.current_stage = 3
                    update_stage_card(3, 'running', stage_names[3])
                elif "STEP 4:" in text or "VERIFYING STYLED OUTPUT" in text:
                    if self.current_stage == 3:
                        update_stage_card(3, 'complete', stage_names[3])
                    self.current_stage = 4
                    update_stage_card(4, 'running', stage_names[4])
                elif "STEP 5:" in text or "CONVERTING TO APA" in text:
                    if self.current_stage == 4:
                        update_stage_card(4, 'complete', stage_names[4])
                    self.current_stage = 5
                    update_stage_card(5, 'running', stage_names[5])
                elif "STEP 6:" in text or "PLAGIARISM DETECTION" in text:
                    if self.current_stage == 5:
                        update_stage_card(5, 'complete', stage_names[5])
                    self.current_stage = 6
                    update_stage_card(6, 'running', stage_names[6])
                elif "STEP 7:" in text or "POLICY ALIGNMENT" in text:
                    if self.current_stage == 6:
                        update_stage_card(6, 'complete', stage_names[6])
                    self.current_stage = 7
                    update_stage_card(7, 'running', stage_names[7])
                
                # Track iteration numbers for Stage 1
                if self.current_stage == 1:
                    # Match "ITERATION 1/3" format (all caps)
                    if re.search(r'ITERATION (\d+)', text):
                        iter_match = re.search(r'ITERATION (\d+)', text)
                        new_iteration = int(iter_match.group(1))
                        if new_iteration != self.current_iteration:
                            self.current_iteration = new_iteration
                            if self.current_iteration not in self.iteration_data:
                                self.iteration_data[self.current_iteration] = {
                                    'links': [],
                                    'topics': [],
                                    'status': [],
                                    'critique': '',
                                    'gaps': []
                                }
                            # Don't update card yet - wait for some content
                    
                    # Extract iteration details
                    if self.current_iteration > 0:
                        iter_info = self.iteration_data[self.current_iteration]
                        
                        # Extract links - more comprehensive pattern (deduplicate)
                        if 'http' in text.lower():
                            urls = re.findall(r'https?://[^\s<>"\']+', text)
                            for url in urls:
                                if url not in iter_info['links']:
                                    iter_info['links'].append(url)
                        
                        # Extract topics
                        if 'Researching' in text or 'Processing topic' in text or 'Deep research' in text:
                            topic = text.split(':')[-1].strip() if ':' in text else text.strip()
                            if topic and len(topic) < 200:
                                if topic not in iter_info['topics']:
                                    iter_info['topics'].append(topic)
                        
                        # Extract status updates - be more aggressive
                        if any(marker in text for marker in ['✓', '✗', 'Generated', 'Complete', 'Fetching', 'Processing', 'Extracted', 'added', 'evidence']):
                            clean_text = text.strip()
                            if clean_text and clean_text not in iter_info['status'] and len(clean_text) < 300:
                                iter_info['status'].append(clean_text)
                        
                        # Only update card when we have enough content for this NEW iteration
                        if self.current_iteration > self.last_rendered_iteration:
                            if len(iter_info['links']) >= 1 or len(iter_info['status']) >= 2:
                                self.update_stage_1_card('running')
                
                # Extract metrics and update other stages
                if self.current_stage > 1:
                    self.extract_and_display_metrics(text)
            
            def update_stage_1_card(self, status):
                """Custom update for Stage 1 with detailed iteration info using st.status()"""
                if status == 'complete':
                    stage_1_status.update(label="✅ Stage 1: Iterative Refinement - Complete", state="complete", expanded=False)
                else:
                    # Only render NEW iterations to avoid duplicates
                    if self.current_iteration > self.last_rendered_iteration:
                        with stage_1_status:
                            iter_info = self.iteration_data[self.current_iteration]
                            
                            st.markdown(f"**📌 Iteration {self.current_iteration}**")
                            
                            # Show links being researched
                            unique_links = list(set(iter_info['links']))[:5]
                            if unique_links:
                                st.markdown("🔗 **Researching:**")
                                for link in unique_links:
                                    display_link = link[:80] + '...' if len(link) > 80 else link
                                    st.code(display_link, language=None)
                            
                            # Show recent status updates
                            if iter_info['status']:
                                recent_status = [s for s in iter_info['status'][-3:] if s and len(s) < 200]
                                for status_line in recent_status:
                                    st.text(status_line[:120])
                            
                            st.divider()
                        
                        self.last_rendered_iteration = self.current_iteration
            
            def extract_and_display_metrics(self, text):
                # Skip debug and info noise
                if "[DEBUG]" in text or "[INFO]" in text:
                    return
                
                # Get existing messages and append new ones
                content_lines = stage_messages.get(self.current_stage, []).copy()
                metrics = {}
                matched_content = False
                
                # Stage 2: Writing Style Retrieval
                if self.current_stage == 2:
                    if "Retrieved style:" in text:
                        content_lines.append(text.strip())
                        matched_content = True
                    if "Speaker:" in text:
                        content_lines.append(text.strip())
                        matched_content = True
                    if "Audience:" in text:
                        content_lines.append(text.strip())
                        matched_content = True
                    if "No style found" in text:
                        content_lines.append(text.strip())
                        matched_content = True
                
                # Stage 3: Styled Output Generation
                if self.current_stage == 3:
                    if "Styled output generated successfully" in text:
                        content_lines.append("✓ Styled output generated successfully")
                        matched_content = True
                    if "Style applied:" in text:
                        content_lines.append(text.strip())
                        matched_content = True
                    if "Output length:" in text and "characters" in text:
                        if match := re.search(r'(\d+) characters', text):
                            metrics['Output Length'] = f"{match.group(1)} chars"
                            matched_content = True
                    if "Citations:" in text and "instances" in text:
                        if match := re.search(r'(\d+) instances', text):
                            metrics['Citations'] = match.group(1)
                            matched_content = True
                    if "All citations valid" in text:
                        content_lines.append("✓ All citations valid")
                        matched_content = True
                    if "Invalid citations:" in text:
                        content_lines.append(text.strip())
                        matched_content = True
                
                # Stage 4: Citation Verification
                if self.current_stage == 4:
                    if "Citation verification complete" in text or "Verifying citations" in text:
                        content_lines.append(text.strip())
                        matched_content = True
                    if "citations found" in text.lower() and "verified" not in text.lower():
                        if match := re.search(r'(\d+) citations', text):
                            metrics['Total Citations'] = match.group(1)
                            matched_content = True
                    if "verified successfully" in text.lower():
                        content_lines.append("✓ " + text.strip())
                        matched_content = True
                    if "Verification rate" in text:
                        content_lines.append(text.strip())
                        matched_content = True
                    if "AUTO-FIX" in text.upper():
                        content_lines.append("🔧 Auto-fix applied")
                        matched_content = True
                
                # Stage 5: APA Conversion
                if self.current_stage == 5:
                    if "APA conversion complete" in text or "Converting to APA" in text:
                        content_lines.append(text.strip())
                        matched_content = True
                    if "Citations converted" in text:
                        if match := re.search(r'(\d+)', text):
                            metrics['Converted'] = match.group(1)
                            matched_content = True
                    if "References generated" in text or "References section" in text:
                        content_lines.append("✓ " + text.strip())
                        matched_content = True
                
                # Stage 6: Plagiarism Detection
                if self.current_stage == 6:
                    if "Plagiarism check complete" in text:
                        content_lines.append(text.strip())
                        matched_content = True
                    if "Overall risk:" in text:
                        if match := re.search(r'Overall risk: (\w+)', text):
                            metrics['Risk Level'] = match.group(1).upper()
                            content_lines.append(f"Risk Level: {match.group(1).upper()}")
                            matched_content = True
                    if "HIGH" in text.upper() and "risk" in text.lower():
                        content_lines.append("⚠️ " + text.strip())
                        matched_content = True
                    if "AUTO-FIX" in text.upper():
                        content_lines.append("🔧 Auto-fix applied")
                        matched_content = True
                
                # Stage 7: Policy Check
                if self.current_stage == 7:
                    if "Policy check:" in text or "Policy alignment" in text:
                        content_lines.append(text.strip())
                        matched_content = True
                    if "COMPLIANT" in text.upper() or "APPROVED" in text.upper():
                        content_lines.append("✓ " + text.strip())
                        matched_content = True
                    if "violation" in text.lower() or "non-compliant" in text.lower():
                        content_lines.append("⚠️ " + text.strip())
                        matched_content = True
                    if "AUTO-FIX" in text.upper():
                        content_lines.append("🔧 Auto-fix applied")
                        matched_content = True
                
                # Update card with captured info - only if we matched something
                if matched_content or metrics:
                    # Keep only last 4 content lines for display (clean and focused)
                    recent_content = content_lines[-4:] if len(content_lines) > 4 else content_lines
                    stage_messages[self.current_stage] = content_lines  # Store all
                    
                    # HTML-escape content and format as plain text
                    escaped_lines = [html.escape(line) for line in recent_content]
                    formatted_content = "\n".join(escaped_lines)
                    
                    update_stage_card(
                        self.current_stage,
                        'running',
                        stage_names[self.current_stage],
                        content=formatted_content,
                        metrics=metrics if metrics else None
                    )
        
        progress_capture = ProgressCapture()
        
        # Monkey-patch print to capture output
        original_print = print
        def custom_print(*args, **kwargs):
            text = " ".join(str(arg) for arg in args)
            progress_capture.parse_and_update(text)
            original_print(*args, **kwargs)
        
        # Temporarily replace print
        import builtins
        builtins.print = custom_print
        
        # Run the complete pipeline
        try:
            user_max_iterations = st.session_state.get("max_iterations", 3)
            results = asyncio.run(
                process_with_iterative_refinement_and_style(
                    query=user_query,
                    sources=sources,
                    max_iterations=user_max_iterations,
                    context_details=context_details,
                    style=selected_style if selected_style else None,
                    enable_policy_check=True
                )
            )
            
            # Mark remaining stages as complete and update with final content
            for i in range(1, 8):
                if stage_status[i] == 'running':
                    update_stage_card(i, 'complete', stage_names[i])
            
            # Update Stage 2 with final style info
            if results.get('style_used'):
                style_info = results['style_used']
                with stage_2_status:
                    st.write(f"**Style:** {style_info.get('name', 'Unknown')}")
                    st.write(f"**Speaker:** {style_info.get('speaker', 'Unknown')}")
                    st.write(f"**Audience:** {style_info.get('audience', 'General')}")
            
            # Update Stage 3 with styled output preview
            if results.get('styled_output') and results['styled_output'].get('styled_output'):
                with stage_3_status:
                    styled_output = results['styled_output']['styled_output']
                    
                    # Show key metrics
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        st.metric("Output Length", f"{len(styled_output)} chars")
                    with col2:
                        citations = results['styled_output'].get('citations_found', 0)
                        st.metric("Citations", citations)
                    with col3:
                        evidence = results['styled_output'].get('unique_evidence_cited', 0)
                        st.metric("Evidence IDs", evidence)
                    
                    st.divider()
                    
                    # Show expandable full output
                    with st.expander("📄 View Generated Styled Output", expanded=True):
                        st.text_area("Styled Output:", styled_output, height=300, disabled=True, label_visibility="collapsed")
            
            # Restore original print
            builtins.print = original_print
            
            # Store results in session state
            st.session_state["pipeline_results"] = results
            
            st.markdown("---")
            st.markdown("### ✅ Pipeline Completed Successfully!")
            
            # === ARCHITECTURAL IMPROVEMENTS VALIDATION (like test_complete_pipeline.py) ===
            st.markdown("#### 🔧 Architectural Improvements Validation")
            
            # Import validation functions from test module
            import sys
            sys.path.insert(0, '/workspaces/pluma-writer')
            from test_complete_pipeline import (
                validate_evidence_stability,
                validate_claim_outline,
                validate_citation_discipline,
                validate_sentence_level_verification,
                validate_boilerplate_filtering,
                validate_policy_compliance_enum
            )
            
            # Run validations
            evidence_stability = validate_evidence_stability(results.get('iterations', []))
            claim_validation = validate_claim_outline(results)
            
            styled_result = results.get('styled_output', {})
            citation_discipline = {}
            if styled_result.get('success'):
                styled_text = styled_result.get('styled_output', '')
                citation_discipline = validate_citation_discipline(styled_text)
            
            verification_result = results.get('citation_verification', {})
            sentence_verify = validate_sentence_level_verification(verification_result)
            
            plagiarism_result = results.get('plagiarism_analysis', {})
            boilerplate = validate_boilerplate_filtering(plagiarism_result)
            
            policy_result = results.get('policy_check', {})
            policy_enum = validate_policy_compliance_enum(policy_result)
            
            # Calculate score
            score = 0
            max_score = 6
            
            if evidence_stability.get('is_cumulative'):
                score += 1
            if claim_validation.get('claim_outline_used'):
                score += 1
            if citation_discipline.get('disciplined'):
                score += 1
            if sentence_verify.get('is_sentence_level'):
                score += 1
            if boilerplate.get('filtering_active'):
                score += 1
            if policy_enum.get('is_strict_enum') and not policy_enum.get('is_unknown'):
                score += 1
            
            grade_percent = (score / max_score) * 100
            grade_letter = 'A+' if grade_percent >= 95 else 'A' if grade_percent >= 90 else 'B+' if grade_percent >= 85 else 'B' if grade_percent >= 80 else 'C+'
            
            # Display improvements in columns
            col1, col2 = st.columns(2)
            
            with col1:
                st.markdown("**1️⃣ Evidence Stability (Cumulative Store)**")
                if evidence_stability.get('enabled'):
                    counts = evidence_stability.get('evidence_counts', [])
                    st.write(f"Evidence counts: {' → '.join(map(str, counts))}")
                    if evidence_stability.get('is_cumulative'):
                        st.success("✅ Evidence is cumulative (no loss across iterations)")
                    else:
                        st.error("❌ Evidence loss detected!")
                
                st.markdown("**2️⃣ Claim-Based Style (Prevents Drift)**")
                if claim_validation.get('claim_outline_used'):
                    st.success(f"✅ Claim outline used: {claim_validation.get('claim_count')} claims")
                    st.write(f"Factual sentences: {claim_validation.get('factual_sentences')}/{claim_validation.get('total_sentences')}")
                else:
                    st.warning("⚠️ Freeform prose used (may cause style drift)")
                
                st.markdown("**3️⃣ Citation Discipline (≤2 IDs/sentence)**")
                if citation_discipline:
                    st.write(f"Avg IDs/sentence: {citation_discipline.get('avg_ids_per_sentence', 0):.2f}")
                    st.write(f"Max IDs/sentence: {citation_discipline.get('max_ids_per_sentence', 0)}")
                    if citation_discipline.get('disciplined'):
                        st.success("✅ Citation discipline maintained")
                    else:
                        violations = len(citation_discipline.get('citation_violations', []))
                        st.warning(f"⚠️ {violations} citation violations detected")
            
            with col2:
                st.markdown("**4️⃣ Sentence-Level Verification**")
                if sentence_verify.get('is_sentence_level'):
                    st.success("✅ Sentence-level verification active")
                else:
                    st.warning("⚠️ Segment-level verification (less precise)")
                st.write(f"Verified: {sentence_verify.get('verified_count')}/{sentence_verify.get('total_segments')}")
                st.write(f"Rate: {sentence_verify.get('verification_rate')}")
                
                st.markdown("**5️⃣ Boilerplate Filtering**")
                if boilerplate.get('filtering_active'):
                    st.success("✅ Boilerplate filtering active")
                    st.write(f"Filtered: {boilerplate.get('boilerplate_patterns_filtered')} chunks")
                else:
                    st.warning("⚠️ No boilerplate filtering detected")
                
                st.markdown("**6️⃣ Strict Policy Compliance Enum**")
                if policy_enum.get('is_strict_enum') and not policy_enum.get('is_unknown'):
                    st.success(f"✅ Strict enum: {policy_enum.get('overall_compliance')}")
                    st.write(f"Score: {policy_enum.get('compliance_score', 0):.1%}")
                elif policy_enum.get('is_unknown'):
                    st.error("❌ UNKNOWN status detected")
                else:
                    st.warning(f"⚠️ Non-standard: {policy_enum.get('overall_compliance')}")
            
            # Overall grade
            st.markdown("---")
            st.markdown(f"### 📊 Overall Grade: **{score}/{max_score}** ({grade_percent:.0f}%) - **{grade_letter}**")
            
            progress_bar_html = f"""
            <div style="width: 100%; background: #e0e0e0; border-radius: 10px; height: 30px; margin: 10px 0;">
                <div style="width: {grade_percent}%; background: {'#198754' if grade_percent >= 80 else '#ffc107' if grade_percent >= 60 else '#dc3545'}; 
                            border-radius: 10px; height: 30px; display: flex; align-items: center; justify-content: center; color: white; font-weight: bold;">
                    {grade_percent:.0f}%
                </div>
            </div>
            """
            st.markdown(progress_bar_html, unsafe_allow_html=True)
            
            st.markdown("---")
            
            # Show summary metrics
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                st.metric("Iterations", results.get("total_iterations", 0))
            with col2:
                st.metric("Evidence Count", results.get("final_evidence_count", 0))
            with col3:
                verification = results.get("citation_verification", {})
                verify_rate = verification.get("verification_rate", "0%")
                st.metric("Citation Verified", verify_rate)
            with col4:
                plagiarism = results.get("plagiarism_analysis", {})
                risk_level = plagiarism.get("overall_risk_level", "unknown").upper()
                risk_color = "🟢" if risk_level in ["NONE", "LOW"] else "🟡" if risk_level == "MEDIUM" else "🔴"
                st.metric("Plagiarism Risk", f"{risk_color} {risk_level}")
            
            # === DETAILED ITERATION BREAKDOWN ===
            st.markdown("---")
            st.markdown("### 🔄 Iteration Breakdown")
            
            with st.expander("📊 Per-Iteration Evidence Collection", expanded=False):
                import re
                for i, iteration in enumerate(results.get('iterations', []), 1):
                    st.markdown(f"**Iteration {i}:**")
                    col_a, col_b, col_c = st.columns(3)
                    with col_a:
                        st.metric("Evidence Items", iteration.get('cumulative_evidence_count', 0))
                    with col_b:
                        st.metric("New Evidence", iteration.get('new_evidence_count', 0))
                    with col_c:
                        has_critique = iteration.get('critique') is not None
                        st.metric("Has Critique", "✓" if has_critique else "✗")
                    
                    # Check evidence format
                    evidence_store = iteration.get('results', {}).get('cumulative_evidence_store', [])
                    if evidence_store:
                        first_evidence = evidence_store[0]
                        is_atomic = 'id' in first_evidence and 'claim' in first_evidence
                        last_evidence = evidence_store[-1]
                        
                        st.write(f"Format: {'✓ ATOMIC' if is_atomic else '✗ OLD PARAGRAPH'}")
                        if is_atomic:
                            st.write(f"ID Range: `{first_evidence.get('id')}` → `{last_evidence.get('id')}`")
                            with st.expander(f"Sample Evidence ({first_evidence.get('id')})", expanded=False):
                                st.json(first_evidence)
                        
                        # Check citations in summary
                        summary = iteration.get('results', {}).get('generated_summary', {}).get('summary', '')
                        citations = re.findall(r'\[E\d+(?:,E\d+)*\]', summary)
                        if citations:
                            cited_ids = set()
                            for citation in citations:
                                ids = re.findall(r'E\d+', citation)
                                cited_ids.update(ids)
                            st.write(f"Citations: {len(citations)} instances, {len(cited_ids)} unique IDs")
                            st.write(f"Sample: {', '.join(citations[:5])}")
                        else:
                            st.warning("⚠️ No [ENN] citations found in summary")
                    
                    st.markdown("---")
            
            # === STYLED OUTPUT PREVIEW ===
            st.markdown("### 📄 Styled Output Details")
            
            style_info = results.get('style_used', {})
            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric("Style Name", style_info.get('name', 'None'))
            with col2:
                st.metric("Speaker", style_info.get('speaker', 'None'))
            with col3:
                st.metric("Audience", style_info.get('audience', 'None'))
            
            if styled_result.get('success'):
                styled_text = styled_result.get('styled_output', '')
                
                col1, col2, col3 = st.columns(3)
                with col1:
                    st.metric("Output Length", f"{len(styled_text)} chars")
                with col2:
                    st.metric("Model", styled_result.get('model_used', 'Unknown'))
                with col3:
                    citations_found = styled_result.get('citations_found', 0)
                    st.metric("Citations", citations_found)
                
                # Citation stats
                if styled_result.get('citations_found') is not None:
                    col1, col2 = st.columns(2)
                    with col1:
                        st.write(f"**Evidence cited:** {styled_result.get('unique_evidence_cited')} unique IDs")
                        st.write(f"**Coverage:** {styled_result.get('citation_coverage')}")
                    with col2:
                        invalid = styled_result.get('invalid_citations')
                        if invalid:
                            st.error(f"⚠️ Invalid citations: {invalid}")
                        else:
                            st.success("✓ All citations valid")
                
                # Preview
                with st.expander("📖 Styled Output Preview (first 1000 chars)", expanded=False):
                    st.text_area("Preview", styled_text[:1000] + ("\n... (truncated)" if len(styled_text) > 1000 else ""), height=300, disabled=True)
            
            st.markdown("---")
            
            # Auto-fix metrics
            st.markdown("#### 🔧 Auto-Fix Summary")
            col1, col2, col3 = st.columns(3)
            
            # Citation fixes
            citation_fix = results.get("citation_fix_result", {})
            with col1:
                fixes = citation_fix.get("fixes_applied", 0)
                if fixes > 0:
                    st.metric("Citation Fixes", fixes, delta=f"+{fixes} applied", delta_color="normal")
                else:
                    st.metric("Citation Fixes", "None needed", delta="✓")
            
            # Plagiarism fixes
            plagiarism_fix = results.get("plagiarism_fix_result", {})
            with col2:
                fixes = plagiarism_fix.get("fixes_applied", 0)
                if fixes > 0:
                    st.metric("Plagiarism Fixes", fixes, delta=f"+{fixes} applied", delta_color="normal")
                else:
                    st.metric("Plagiarism Fixes", "None needed", delta="✓")
            
            # Policy fixes
            policy_fix = results.get("policy_fix_result", {})
            with col3:
                fixes = policy_fix.get("fixes_applied", 0)
                if fixes > 0:
                    st.metric("Policy Fixes", fixes, delta=f"+{fixes} applied", delta_color="normal")
                else:
                    st.metric("Policy Fixes", "None needed", delta="✓")
            
            # === DETAILED VERIFICATION RESULTS ===
            st.markdown("---")
            st.markdown("### 🔍 Verification & Compliance Details")
            
            tab_verify, tab_plagia, tab_policy = st.tabs(["Citation Verification", "Plagiarism Analysis", "Policy Compliance"])
            
            with tab_verify:
                verification = results.get("citation_verification", {})
                if verification.get("success"):
                    col1, col2 = st.columns(2)
                    with col1:
                        st.metric("Total Segments", verification.get("total_segments", 0))
                        st.metric("Verified", verification.get("verified_segments", 0))
                    with col2:
                        st.metric("Unverified", verification.get("unverified_segments", 0))
                        rate = verification.get("verification_rate", "0%")
                        rate_float = float(rate.strip('%'))
                        st.metric("Verification Rate", rate, delta="✓ Excellent" if rate_float >= 95 else "⚠️ Needs improvement")
                    
                    # Show unverified segments
                    unverified_segments = [s for s in verification.get("segments", []) if s.get("verified") == "No"]
                    if unverified_segments:
                        st.warning(f"**{len(unverified_segments)} Unverified Segments:**")
                        for seg in unverified_segments[:5]:
                            with st.expander(f"Segment {seg.get('segment_number')}", expanded=False):
                                st.write(f"**Sentence:** {seg.get('sentence', '')}")
                                st.write(f"**Reason:** {seg.get('verification_reason', 'Unknown')}")
                                st.write(f"**Citations:** {seg.get('citations', [])}")
                    else:
                        st.success("✅ All citations verified successfully!")
                else:
                    st.warning("Citation verification not available")
            
            with tab_plagia:
                plagiarism = results.get("plagiarism_analysis", {})
                if plagiarism.get("success"):
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        risk = plagiarism.get("overall_risk_level", "unknown").upper()
                        risk_icon = "🟢" if risk in ["NONE", "LOW"] else "🟡" if risk == "MEDIUM" else "🔴"
                        st.metric("Overall Risk", f"{risk_icon} {risk}")
                    with col2:
                        stats = plagiarism.get("statistics", {})
                        st.metric("Chunks Analyzed", stats.get("total_chunks", 0))
                        st.metric("High Risk", stats.get("high_risk_count", 0))
                    with col3:
                        st.metric("Sources Found", len(plagiarism.get("sources_found", [])))
                        st.metric("Boilerplate Filtered", stats.get("boilerplate_filtered", 0))
                    
                    # Show sources
                    sources = plagiarism.get("sources_found", [])
                    if sources:
                        st.markdown("**Sources with Matches:**")
                        for i, src in enumerate(sources[:5], 1):
                            st.write(f"{i}. {src.get('title', 'Unknown')} ({src.get('match_count', 0)} matches)")
                    
                    # Show high-risk chunks
                    chunks = plagiarism.get("chunks", [])
                    high_risk = [c for c in chunks if c.get("risk_level") in ["HIGH", "CRITICAL"]]
                    if high_risk:
                        st.warning(f"**{len(high_risk)} High-Risk Chunks:**")
                        for chunk in high_risk[:3]:
                            with st.expander(f"Risk Score: {chunk.get('overall_risk_score', 0):.2f}", expanded=False):
                                st.write(f"**Text:** {chunk.get('text', '')[:200]}...")
                                st.write(f"**Risk Level:** {chunk.get('risk_level')}")
                else:
                    st.warning("Plagiarism analysis not available")
            
            with tab_policy:
                policy = results.get("policy_check", {})
                if policy.get("success"):
                    col1, col2 = st.columns(2)
                    with col1:
                        compliance = policy.get("overall_compliance", "UNKNOWN")
                        icon = "✅" if compliance == "COMPLIANT" else "⚠️" if compliance == "NEEDS_REVISION" else "❌"
                        st.metric("Compliance Status", f"{icon} {compliance}")
                    with col2:
                        score = policy.get("compliance_score", 0.0)
                        st.metric("Compliance Score", f"{score:.1%}")
                    
                    # Show violations
                    violations = policy.get("violations", [])
                    if violations:
                        st.error(f"**{len(violations)} Policy Violations:**")
                        for viol in violations[:5]:
                            with st.expander(f"{viol.get('violation_type', 'Unknown')}", expanded=False):
                                st.write(f"**Description:** {viol.get('description', '')}")
                                st.write(f"**Location:** {viol.get('problematic_text', '')[:200]}")
                                st.write(f"**Severity:** {viol.get('severity', 'UNKNOWN')}")
                                if viol.get('suggested_fix'):
                                    st.info(f"**Suggested Fix:** {viol.get('suggested_fix')}")
                    else:
                        st.success("✅ No policy violations detected!")
                else:
                    st.warning("Policy check not available")
            
            st.markdown("---")
            st.markdown("### 📝 Generated Speech")
            
            # Tabs for different outputs
            tab1, tab2, tab3 = st.tabs(["Styled Output", "APA Format", "Raw Summary"])
            
            with tab1:
                styled_result = results.get("styled_output", {})
                if styled_result.get("success"):
                    styled_text = styled_result.get("styled_output", "")
                    st.text_area("Styled Speech", styled_text, height=400, key="styled_display")
                    
                    # Download buttons
                    ts = datetime.now().strftime("%Y%m%d-%H%M%S")
                    speaker = st.session_state.get("selected_speaker_input", "Speaker")
                    base_name = f"speech_{speaker}_{ts}"
                    
                    col1, col2 = st.columns(2)
                    with col1:
                        docx_bytes = make_docx_bytes(styled_text, title=f"Speech - {speaker}")
                        st.download_button(
                            "⬇️ Download Styled Output (DOCX)",
                            data=docx_bytes,
                            file_name=f"{base_name}_styled.docx",
                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                            use_container_width=True
                        )
                    with col2:
                        pdf_bytes = make_pdf_bytes(styled_text, title=f"Speech - {speaker}")
                        st.download_button(
                            "⬇️ Download Styled Output (PDF)",
                            data=pdf_bytes,
                            file_name=f"{base_name}_styled.pdf",
                            mime="application/pdf",
                            use_container_width=True
                        )
                else:
                    st.error(f"❌ Styled output failed: {styled_result.get('error', 'Unknown error')}")
            
            with tab2:
                apa_result = results.get("styled_output_apa", {})
                if apa_result.get("success"):
                    apa_text = apa_result.get("apa_output", "")
                    st.text_area("APA Format Speech", apa_text, height=400, key="apa_display")
                    
                    # Download buttons
                    ts = datetime.now().strftime("%Y%m%d-%H%M%S")
                    speaker = st.session_state.get("selected_speaker_input", "Speaker")
                    base_name = f"speech_{speaker}_{ts}"
                    
                    col1, col2 = st.columns(2)
                    with col1:
                        docx_bytes = make_docx_bytes(apa_text, title=f"Speech (APA) - {speaker}")
                        st.download_button(
                            "⬇️ Download APA Format (DOCX)",
                            data=docx_bytes,
                            file_name=f"{base_name}_apa.docx",
                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                            use_container_width=True
                        )
                    with col2:
                        pdf_bytes = make_pdf_bytes(apa_text, title=f"Speech (APA) - {speaker}")
                        st.download_button(
                            "⬇️ Download APA Format (PDF)",
                            data=pdf_bytes,
                            file_name=f"{base_name}_apa.pdf",
                            mime="application/pdf",
                            use_container_width=True
                        )
                else:
                    st.warning("⚠️ APA format conversion not available")
            
            with tab3:
                final_summary = results.get("final_summary", {}).get("summary", "")
                st.text_area("Raw Summary (Before Styling)", final_summary, height=400, key="summary_display")
            
            # Show detailed reports in expanders
            st.markdown("---")
            st.markdown("### 📊 Detailed Reports")
            
            with st.expander("🔍 Citation Verification Report", expanded=False):
                verification = results.get("citation_verification", {})
                if verification.get("success"):
                    st.json(verification)
                else:
                    st.warning("Citation verification not available")
            
            with st.expander("📋 Plagiarism Analysis Report", expanded=False):
                plagiarism = results.get("plagiarism_analysis", {})
                if plagiarism.get("success"):
                    st.json(plagiarism)
                else:
                    st.warning("Plagiarism analysis not available")
            
            with st.expander("⚖️ Policy Compliance Report", expanded=False):
                policy = results.get("policy_check", {})
                if policy.get("success"):
                    st.json(policy)
                else:
                    st.warning("Policy check not available")
            
            with st.expander("🔧 Auto-Fix Details", expanded=False):
                st.markdown("**Citation Fixes:**")
                st.json(citation_fix)
                st.markdown("**Plagiarism Fixes:**")
                st.json(plagiarism_fix)
                st.markdown("**Policy Fixes:**")
                st.json(policy_fix)
            
            with st.expander("📦 Complete Results (JSON)", expanded=False):
                st.json(results)
                
                # Download complete JSON
                json_bytes = json.dumps(results, indent=2, ensure_ascii=False).encode('utf-8')
                ts = datetime.now().strftime("%Y%m%d-%H%M%S")
                st.download_button(
                    "⬇️ Download Complete Results (JSON)",
                    data=json_bytes,
                    file_name=f"pipeline_results_{ts}.json",
                    mime="application/json",
                    use_container_width=True
                )
            
        except Exception as e:
            # Restore original print on error
            builtins.print = original_print
            
            # Mark current stage as error
            if progress_capture.current_stage > 0:
                update_stage_card(progress_capture.current_stage, 'error', stage_names[progress_capture.current_stage])
            
            st.error(f"❌ Pipeline failed: {str(e)}")
            import traceback
            with st.expander("🐛 Error Details", expanded=True):
                st.code(traceback.format_exc())